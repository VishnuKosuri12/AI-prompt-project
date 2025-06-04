import os
from pptx import Presentation

def create_presentation():
    # Load the template
    prs = Presentation('tasks/Covesto_Master_130324.pptx')
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "AI Coding"
    subtitle.text = "ChemTrack Project"
    
    # Section 1: AI Coding
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    slide.shapes.title.text = "What is AI Coding?"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Modern AI coding tools provide:"
    p = tf.add_paragraph()
    p.text = "• Problem solving"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  ▪ Cova"
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Coding Assistant"
    p.level = 1
    
    assistants = ["GitHub Copilot", "Amazon Q", "Tabnine", "Codeium", "JetBrains AI Assistant"]
    for assistant in assistants:
        p = tf.add_paragraph()
        p.text = f"  ▪ {assistant}"
        p.level = 2
        
    p = tf.add_paragraph()
    p.text = "• Code authoring"
    p.level = 1
    
    authors = ["Cline, VS Code, AWS Bedrock"]
    for author in authors:
        p = tf.add_paragraph()
        p.text = f"  ▪ {author}"
        p.level = 2
    
    # For all remaining slides, use the layout from the third slide of the template (layout #12 - Text)
    
    # Section 2: How I got here
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "How I got here"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    stats = [
        "• Y Combinator -- 25% of their startups are using AI to build their code bases",
        "• AI boosts developer productivity 10-30%",
        "• AI assisted developers write 12-15% more code",
        "• 25% of Google's code is now AI generated"
    ]
    
    for i, stat in enumerate(stats):
        if i == 0:
            tf.text = stat
        else:
            p = tf.add_paragraph()
            p.text = stat
    
    # Section 3: I have questions
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "I have questions"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    questions = [
        "• Can it handle a large code base?",
        "• Can it keep code organized and clean?",
        "• What about modifying a large code base?",
        "• Can it debug code that's broken?"
    ]
    
    for i, question in enumerate(questions):
        if i == 0:
            tf.text = question
        else:
            p = tf.add_paragraph()
            p.text = question
    
    # Section 4: My Goals
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "My Goals"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    goals = [
        "• Build a large codebase using AI",
        "• Learn how to prompt it to get quality results",
        "• Try different approaches to see what it generates",
        "• Understand what it can and can't generate",
        "• Bring my learnings back to the organization"
    ]
    
    for i, goal in enumerate(goals):
        if i == 0:
            tf.text = goal
        else:
            p = tf.add_paragraph()
            p.text = goal
    
    # Section 5: What are we building?
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "What are we building?"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    tf.text = "• ChemTrack -- a chemical tracking application"
    
    # Architecture slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "Architecture"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    architecture = [
        "• Internally facing application",
        "• Load balancers for interfacing",
        "• ECS environment running multiple containers",
        "• RDS PostgreSQL database"
    ]
    
    for i, arch in enumerate(architecture):
        if i == 0:
            tf.text = arch
        else:
            p = tf.add_paragraph()
            p.text = arch
    
    # Technologies slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "Technologies used"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    technologies = [
        "• Python: Flask, Jinja, FastAPI",
        "• HTML, Javascript",
        "• SQL (PostgreSQL dialect)",
        "• Docker",
        "• Bash -- scripting",
        "• AWS: environment, CLI, API (boto3), CloudFormation"
    ]
    
    for i, tech in enumerate(technologies):
        if i == 0:
            tf.text = tech
        else:
            p = tf.add_paragraph()
            p.text = tech
    
    # Section 6: What did we learn?
    # Prompting slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "What did we learn? - Prompting"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    prompting = [
        "• Detail is good but not always necessary or beneficial",
        "• Remind it where things are located and how you want things structured",
        "• Be specific if you are looking for specific outcomes",
        "• Use shared files for common information"
    ]
    
    for i, prompt in enumerate(prompting):
        if i == 0:
            tf.text = prompt
        else:
            p = tf.add_paragraph()
            p.text = prompt
    
    # Generated Code slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "What did we learn? - Generated Code"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    gen_code = [
        "• Structure and organization",
        "• Code appearance",
        "• Maintainability"
    ]
    
    for i, code in enumerate(gen_code):
        if i == 0:
            tf.text = code
        else:
            p = tf.add_paragraph()
            p.text = code
    
    # What it can do slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "What did we learn? - What it can do"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    capabilities = [
        "• Add new modules",
        "• Feature enhancements",
        "• Database enhancements",
        "• Generate database data",
        "• Diagram the application",
        "• Refactor code structure",
        "• Refactor user interface"
    ]
    
    for i, capability in enumerate(capabilities):
        if i == 0:
            tf.text = capability
        else:
            p = tf.add_paragraph()
            p.text = capability
    
    # Section 7: Working with AI
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "Working with AI - Commands and Responses"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    tf.text = "Seasoned developer and the gorilla with a hammer"
    p = tf.add_paragraph()
    p.text = "(image placeholder)"
    
    # Good things slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "Working with AI - Good things"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    good_things = [
        "• Maintaining documentation (readme)",
        "• Building and maintaining scripts for build and deploy",
        "• Using interesting technologies and/or techniques",
        "• Showing signs of good application structure",
        "• Refactoring code",
        "• Pretty good at layout"
    ]
    
    for i, good in enumerate(good_things):
        if i == 0:
            tf.text = good
        else:
            p = tf.add_paragraph()
            p.text = good
    
    # Odd things slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "Working with AI - Odd things"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    odd_things = [
        "• Flipping between two solutions when debugging",
        "• Proposing multiple solutions but picking one without asking",
        "• Diff Edit Mismatch",
        "  (image placeholder)",
        "• Bad code -- that it caught itself",
        "• Running out of steam and writing incomplete files -- then fixing it"
    ]
    
    for i, odd in enumerate(odd_things):
        if i == 0:
            tf.text = odd
        else:
            p = tf.add_paragraph()
            p.text = odd
    
    # Bad things slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "Working with AI - Bad things"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    bad_things = [
        "• Removing my changes",
        "• Making a change in some modules but not in others",
        "• Moving files from one place to another (duplicating)",
        "• Bad code (occasional)",
        "• Not great at debugging"
    ]
    
    for i, bad in enumerate(bad_things):
        if i == 0:
            tf.text = bad
        else:
            p = tf.add_paragraph()
            p.text = bad
    
    # Things to watch slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "Working with AI - Things to watch"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    watch_things = [
        "• Be clear about names and paths",
        "• Be clear about technolgies you want used",
        "• Remind it of changes that will impact many parts of the application"
    ]
    
    for i, watch in enumerate(watch_things):
        if i == 0:
            tf.text = watch
        else:
            p = tf.add_paragraph()
            p.text = watch
    
    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    slide.shapes.title.text = "Summary"
    
    # Find body placeholder by its type
    body = None
    for shape in slide.shapes:
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 2:  # BODY
            body = shape
            break
    
    tf = body.text_frame
    tf.text = "• AI coding tools can significantly enhance developer productivity"
    p = tf.add_paragraph()
    p.text = "• With proper prompting and oversight, AI can build large, maintainable codebases"
    p = tf.add_paragraph()
    p.text = "• Understanding AI's strengths and limitations is key to effective collaboration"
    p = tf.add_paragraph()
    p.text = "• ChemTrack project demonstrates the real-world application of AI coding"
    
    # Save the presentation
    output_path = 'AI_Coding_Presentation.pptx'
    prs.save(output_path)
    print(f"Presentation saved as {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
